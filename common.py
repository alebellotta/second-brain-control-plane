"""Shared helpers for the local second brain: chunking, Ollama calls, Chroma access,
and document-to-note conversion (PDF/DOCX/PPTX/TXT)."""
import hashlib
import json
import logging
import os
import re
import time
from datetime import datetime, timezone
from pathlib import Path

import requests

import policy

VAULT_DIR = Path(os.environ.get("SECOND_BRAIN_VAULT", str(Path.home() / "Documents" / "SecondBrain")))
HOME_DIR = Path(os.environ.get("SECOND_BRAIN_HOME", str(Path.home() / ".second-brain")))
CHROMA_DIR = HOME_DIR / "chroma"
LOG_DIR = HOME_DIR / "logs"
SOURCES_DIR = VAULT_DIR / "Sources"
NOTES_DIR = VAULT_DIR / "Notes"

POLICY_LOG_PATH = LOG_DIR / "policy_decisions.jsonl"

OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")
EMBED_MODEL = "nomic-embed-text"
TAG_MODEL = "llama3.2"
DIGEST_MODEL = "llama3.2"

# Bumped whenever the extraction logic for a format changes materially enough
# that previously-generated notes might read differently if reprocessed.
# Recorded per-note in the frontmatter (extractor_version) so a manifest scan
# can tell which notes were produced by an older version of the pipeline.
EXTRACTOR_VERSIONS = {
    ".pdf": "pdf-pymupdf-v1",
    ".docx": "docx-python-docx-v1",
    ".pptx": "pptx-python-pptx-v1",
    ".txt": "text-plain-v1",
    ".md": "text-plain-v1",
}
CHUNKING_VERSION = "paragraph-blocks-v1"  # bump if chunk_markdown()'s splitting strategy changes

COLLECTION_NAME = "notes"

# Vault folders the watcher should ignore for direct note indexing (Obsidian
# config, our own generated output, and the raw documents in Sources/ that get
# converted into separate notes rather than indexed directly).
IGNORED_DIR_NAMES = {".obsidian", "_Suggestions", "Reviews", "Sources"}

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100


def setup_logging(name: str) -> logging.Logger:
    LOG_DIR.mkdir(parents=True, exist_ok=True)
    logger = logging.getLogger(name)
    logger.setLevel(logging.INFO)
    if not logger.handlers:
        handler = logging.FileHandler(LOG_DIR / f"{name}.log")
        handler.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s"))
        logger.addHandler(handler)
        logger.addHandler(logging.StreamHandler())
    return logger


def is_ignored_path(path: Path) -> bool:
    try:
        rel = path.relative_to(VAULT_DIR)
    except ValueError:
        return True
    return any(part in IGNORED_DIR_NAMES for part in rel.parts[:-1]) or rel.parts[0] in IGNORED_DIR_NAMES


def relpath(path: Path) -> str:
    return str(path.relative_to(VAULT_DIR))


def chunk_markdown(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split on headings (##/###) first, then into overlapping blocks on paragraph boundaries."""
    sections = re.split(r"(?=^#{1,3} .+$)", text, flags=re.MULTILINE)
    chunks: list[str] = []
    for section in sections:
        section = section.strip()
        if not section:
            continue
        if len(section) <= chunk_size:
            chunks.append(section)
            continue
        paragraphs = section.split("\n\n")
        current = ""
        for para in paragraphs:
            if len(current) + len(para) + 2 <= chunk_size:
                current = f"{current}\n\n{para}" if current else para
            else:
                if current:
                    chunks.append(current)
                current = (current[-overlap:] + "\n\n" + para) if current else para
        if current:
            chunks.append(current)
    if not chunks and text.strip():
        chunks.append(text.strip())
    return chunks


def ollama_embed(text: str) -> list[float] | None:
    try:
        resp = requests.post(
            f"{OLLAMA_URL}/api/embeddings",
            json={"model": EMBED_MODEL, "prompt": text},
            timeout=60,
        )
        resp.raise_for_status()
        return resp.json()["embedding"]
    except Exception:
        logging.getLogger("common").exception("Ollama embedding error")
        return None


def system_user_messages(system: str, user: str) -> list[dict]:
    """Convenience for the common case: a single system instruction plus one
    block of (often untrusted) content. Callers that need real multi-turn
    history (chat.py) build their own messages list instead of using this."""
    return [{"role": "system", "content": system}, {"role": "user", "content": user}]


def ollama_chat(model: str, messages: list[dict]) -> str | None:
    """Calls Ollama's /api/chat (not /api/generate): messages carry explicit
    roles instead of being concatenated into one prompt string. This gives
    the model an explicit instruction-vs-data boundary — untrusted note
    content always goes in a `user` message, never mixed into the `system`
    instructions — which most instruct models are trained to weight
    differently. This is a mitigation, not a guarantee: see the paper's
    findings on prompt injection for what this does and doesn't stop."""
    try:
        resp = requests.post(
            f"{OLLAMA_URL}/api/chat",
            json={"model": model, "messages": messages, "stream": False},
            timeout=180,
        )
        resp.raise_for_status()
        return resp.json()["message"]["content"].strip()
    except Exception:
        logging.getLogger("common").exception("Ollama chat error")
        return None


def ollama_chat_json(model: str, messages: list[dict], schema: dict) -> dict | None:
    """Like ollama_chat, but constrains the output to a JSON schema (natively
    supported by Ollama via the "format" parameter). More robust than
    free-text parsing for structured data: the model only has to choose the
    values, not also guess the exact syntax to follow — this removes a whole
    class of parsing failures without needing a more elaborate prompt. Note
    that the schema constrains the response's SHAPE, not the CONTENT of each
    string field — see watcher.py's _sanitize_tag() for why that distinction
    matters in practice."""
    try:
        resp = requests.post(
            f"{OLLAMA_URL}/api/chat",
            json={"model": model, "messages": messages, "format": schema, "stream": False},
            timeout=180,
        )
        resp.raise_for_status()
        return json.loads(resp.json()["message"]["content"])
    except Exception:
        logging.getLogger("common").exception("Ollama structured chat error")
        return None


def _markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    header, *body = rows
    lines = ["| " + " | ".join(header) + " |", "|" + "---|" * len(header)]
    for row in body:
        lines.append("| " + " | ".join(c.replace("\n", " ") for c in row) + " |")
    return "\n".join(lines)


OCR_MIN_TEXT_LENGTH = 20  # below this, a page is almost certainly image-only
OCR_LANGUAGES = "eng"  # e.g. "eng+ita" for multilingual OCR; needs the matching tesseract-lang data


def _ocr_pdf_page(page, dpi: int = 200) -> tuple[str, float | None]:
    """OCR of a rendered PDF page (fallback for scanned PDFs with no text
    layer). Requires the tesseract binary to be installed (e.g.
    brew install tesseract tesseract-lang, or apt install tesseract-ocr on
    Linux); if missing, fails silently and the page is simply left without
    text. Returns (text, mean_word_confidence_0_to_100_or_None) — the
    confidence is surfaced in the note's manifest (ocr_confidence) so a low-
    confidence OCR result is visibly distinguishable from a clean text layer,
    instead of looking identical in the index."""
    try:
        import io

        import pytesseract
        from PIL import Image
    except ImportError:
        return "", None

    pix = page.get_pixmap(dpi=dpi)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    try:
        data = pytesseract.image_to_data(img, lang=OCR_LANGUAGES, output_type=pytesseract.Output.DICT)
    except Exception:
        try:
            data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
        except Exception:
            logging.getLogger("common").exception("OCR failed for page %d", page.number + 1)
            return "", None

    lines: list[list[str]] = []
    confidences: list[float] = []
    current_key = None
    for i, word in enumerate(data["text"]):
        word = word.strip()
        if not word:
            continue
        key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
        if key != current_key:
            lines.append([])
            current_key = key
        lines[-1].append(word)
        conf = float(data["conf"][i])
        if conf >= 0:
            confidences.append(conf)

    text = "\n".join(" ".join(line) for line in lines)
    mean_confidence = sum(confidences) / len(confidences) if confidences else None
    return text, mean_confidence


def _extract_pdf(path: Path) -> tuple[str, dict]:
    """Text via PyMuPDF, one section per page. If a page has no extractable
    text (typical of a scanned PDF), falls back to local OCR before treating
    it as empty. Returns (text, metadata) — metadata includes ocr_confidence
    (average across OCR'd pages, 0-100) only if OCR was actually used."""
    import fitz

    parts: list[str] = []
    ocr_confidences: list[float] = []

    doc = fitz.open(str(path))
    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            text = page.get_text("text").strip()
            if len(text) < OCR_MIN_TEXT_LENGTH:
                ocr_text, ocr_confidence = _ocr_pdf_page(page)
                if len(ocr_text) > len(text):
                    text = ocr_text
                    if ocr_confidence is not None:
                        ocr_confidences.append(ocr_confidence)
            if text:
                parts.append(f"## Page {page_index + 1}\n\n{text}")
    finally:
        doc.close()

    metadata = {}
    if ocr_confidences:
        metadata["ocr_confidence"] = round(sum(ocr_confidences) / len(ocr_confidences), 1)
        logging.getLogger("common").info(
            "OCR used as fallback for %s (mean confidence %.1f)", path, metadata["ocr_confidence"]
        )

    return "\n\n".join(parts), metadata


def _extract_docx(path: Path) -> tuple[str, dict]:
    """Preserves headings and tables as markdown."""
    import docx

    parts: list[str] = []
    document = docx.Document(str(path))

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if style == "title":
            parts.append(f"# {text}")
        elif style.startswith("heading"):
            digits = "".join(c for c in style if c.isdigit())
            level = min(max(int(digits), 1), 6) if digits else 2
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        md = _markdown_table(rows)
        if md:
            parts.append(md)

    return "\n\n".join(parts), {}


def _extract_pptx(path: Path) -> tuple[str, dict]:
    """Preserves title/bullets (with indentation), tables as markdown, and speaker
    notes, one section per slide."""
    from pptx import Presentation

    slides_md: list[str] = []

    prs = Presentation(str(path))
    for slide_index, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_text = title_shape.text_frame.text.strip() if title_shape and title_shape.has_text_frame else ""
        parts = [f"## Slide {slide_index}" + (f": {title_text}" if title_text else "")]

        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(f"{'  ' * para.level}- {text}")
            elif shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                md = _markdown_table(rows)
                if md:
                    parts.append(md)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append("> Speaker notes: " + notes_text.replace("\n", " "))

        slides_md.append("\n\n".join(parts))

    return "\n\n".join(slides_md), {}


def _extract_txt(path: Path) -> tuple[str, dict]:
    return path.read_text(encoding="utf-8", errors="ignore"), {}


WHISPER_MODEL_SIZE = "small"  # tiny/base/small/medium/large-v3: speed vs. accuracy trade-off on CPU
AUDIO_LANGUAGE = None  # None = auto-detect; set e.g. "en" to skip detection and speed things up

_whisper_model = None


def _get_whisper_model():
    """Loads the Whisper model once per process (same pattern as the Chroma
    client): loading takes a few seconds, no point repeating it for every
    audio file in the same watcher session."""
    global _whisper_model
    if _whisper_model is None:
        from faster_whisper import WhisperModel

        _whisper_model = WhisperModel(WHISPER_MODEL_SIZE, device="cpu", compute_type="int8")
    return _whisper_model


def _extract_audio(path: Path) -> tuple[str, dict]:
    """Transcribes an audio recording (e.g. a meeting) with local Whisper
    (faster-whisper, CPU with int8 quantization — no dedicated GPU
    acceleration on this class of hardware, see the quantization note
    elsewhere in this file for llama3.2). Each segment becomes a line with a
    timestamp, useful for navigating long recordings without listening to
    the whole thing again. Returns (text, metadata) with the detected
    language and its confidence, surfaced in the note's manifest the same
    way OCR confidence is for scanned PDFs."""
    model = _get_whisper_model()
    segments, info = model.transcribe(str(path), language=AUDIO_LANGUAGE)

    lines = [
        f"**[{int(s.start) // 60:02d}:{int(s.start) % 60:02d}]** {s.text.strip()}"
        for s in segments
        if s.text.strip()
    ]
    metadata = {
        "audio_language": info.language,
        "audio_language_confidence": round(info.language_probability * 100, 1),
    }
    if not lines:
        return "", metadata

    header = f"Detected language: {info.language} (confidence {info.language_probability:.0%})"
    return header + "\n\n" + "\n\n".join(lines), metadata


AUDIO_EXTENSIONS = {".mp3", ".m4a", ".wav", ".mp4", ".aac", ".flac"}
EXTRACTOR_VERSIONS.update({ext: "audio-faster-whisper-v1" for ext in AUDIO_EXTENSIONS})

SOURCE_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".txt": _extract_txt,
    ".md": _extract_txt,
    **{ext: _extract_audio for ext in AUDIO_EXTENSIONS},
}


def is_source_document(path: Path) -> bool:
    try:
        path.relative_to(SOURCES_DIR)
    except ValueError:
        return False
    return path.suffix.lower() in SOURCE_EXTRACTORS


def note_path_for_source(path: Path) -> Path:
    # Keep the original extension in the name (e.g. "file.pptx.md") instead of
    # replacing it: two sources with the same stem but different extension
    # (e.g. "report.pdf" and "report.pptx") would otherwise generate the same
    # note, overwriting each other.
    rel = path.relative_to(SOURCES_DIR)
    return NOTES_DIR / rel.with_name(rel.name + ".md")


_VERSION_PATTERN = re.compile(r"^(?P<base>.+?)[ _-]v(?P<version>\d+|final)$", re.IGNORECASE)


def _version_rank(stem: str):
    """From 'Report_v03' extracts ('report', (0, 3)); from 'Report_vFINAL' extracts
    ('report', (1, 0)) — FINAL always outranks any number. None if the name has
    no recognizable version marker."""
    match = _VERSION_PATTERN.match(stem.strip())
    if not match:
        return None
    base = match.group("base").strip().lower()
    version = match.group("version").lower()
    rank = (1, 0) if version == "final" else (0, int(version))
    return base, rank


def latest_version_only(paths: list[Path]) -> set[Path]:
    """Among a list of files (typically from the same folder), groups different
    versions of the same document — same base name, same extension, a
    _v01/_v2/_vFINAL... marker — and keeps only the latest per group (vFINAL
    always wins; otherwise the highest number wins). Files without a recognized
    version marker always pass through unchanged. Returns the set of paths to keep."""
    groups: dict[tuple, list[Path]] = {}
    keep: set[Path] = set()
    for p in paths:
        info = _version_rank(p.stem)
        if info is None:
            keep.add(p)
            continue
        base, _ = info
        groups.setdefault((p.parent, base, p.suffix.lower()), []).append(p)
    for group in groups.values():
        best = max(group, key=lambda p: _version_rank(p.stem)[1])
        keep.add(best)
    return keep


_OFFICE_EXTS = {".pptx", ".docx"}


def prefer_pdf_duplicates(paths: list[Path]) -> set[Path]:
    """If the same document exists both as a PDF and as PPTX/DOCX (same base
    name, same folder) — typically a PDF export of the same deck or report —
    keeps only the PDF, which is almost always much smaller for equivalent
    content. Files without a redundant PDF/Office counterpart pass through
    unchanged. Returns the set of paths to keep."""
    groups: dict[tuple, list[Path]] = {}
    for p in paths:
        groups.setdefault((p.parent, p.stem.lower()), []).append(p)

    keep: set[Path] = set()
    for group in groups.values():
        exts = {p.suffix.lower() for p in group}
        if ".pdf" in exts and exts & _OFFICE_EXTS:
            keep.update(p for p in group if p.suffix.lower() == ".pdf")
        else:
            keep.update(group)
    return keep


def _wait_until_stable(path: Path, timeout: float = 180.0, interval: float = 1.5) -> bool:
    """Waits for a file to stop growing/changing before reading it: useful for
    files still being copied (e.g. syncing from a cloud drive) that would
    otherwise be read half-written, producing a corrupt pptx/docx package.
    Returns False if it never stabilizes within the timeout (huge file or a
    stuck copy)."""
    deadline = time.time() + timeout
    try:
        last_size = path.stat().st_size
    except OSError:
        return False
    while time.time() < deadline:
        time.sleep(interval)
        try:
            current_size = path.stat().st_size
        except OSError:
            return False
        if current_size == last_size:
            return True
        last_size = current_size
    return False


def _stored_source_mtime(out_path: Path) -> float | None:
    """Reads the source mtime recorded in the frontmatter at conversion time.
    More robust than comparing the mtimes of two different files on disk:
    cloud-synced folders (OneDrive/iCloud/Dropbox) can "touch" mtimes without
    the content actually changing, which would otherwise cause needless
    reconversions."""
    try:
        text = out_path.read_text(encoding="utf-8")
    except OSError:
        return None
    match = re.search(r'^source_mtime:\s*([0-9.]+)', text, flags=re.MULTILINE)
    return float(match.group(1)) if match else None


def convert_source(path: Path) -> Path | None:
    """Extracts text from a document in Sources/ and writes the corresponding
    markdown note in Notes/ (same folder structure). Returns the path of the
    note written, or None if nothing needed to be (or could be) written.

    Every candidate is run through policy.evaluate() twice: once before
    reading (path/size — cheap, catches the symlink-escape case without ever
    touching file content) and once after extraction (content — secrets/PII
    patterns). Both decisions are appended to the local policy audit log."""
    log = logging.getLogger("common")

    pre_decision = policy.evaluate(path, SOURCES_DIR)
    if pre_decision.action == "deny":
        log.warning("Policy denied %s: %s", path, "; ".join(pre_decision.reasons))
        policy.log_decision(POLICY_LOG_PATH, path, pre_decision)
        return None

    out_path = note_path_for_source(path)
    current_source_mtime = path.stat().st_mtime
    if out_path.exists():
        stored_mtime = _stored_source_mtime(out_path)
        if stored_mtime is not None and stored_mtime >= current_source_mtime:
            return None

    extractor = SOURCE_EXTRACTORS.get(path.suffix.lower())
    if not extractor:
        return None

    if not _wait_until_stable(path):
        log.warning("File still being written/copied after waiting, skipping for now: %s", path)
        return None

    try:
        text, extraction_meta = extractor(path)
    except Exception:
        log.exception("Extraction failed for %s", path)
        return None

    if not text.strip():
        log.warning("No text extracted from %s (scanned PDF with no OCR?)", path)
        return None

    content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]
    decision = policy.evaluate(path, SOURCES_DIR, text=text)
    policy.log_decision(POLICY_LOG_PATH, path, decision, content_hash=content_hash)
    if decision.action == "quarantine":
        log.warning(
            "Note for %s flagged for quarantine (%s): will be written but not embedded "
            "semantically — see policy_decision/sensitivity_flags in its frontmatter",
            path, ", ".join(decision.sensitivity_flags),
        )

    extra_frontmatter = "".join(f"{key}: {json.dumps(value)}\n" for key, value in extraction_meta.items())

    rel = path.relative_to(SOURCES_DIR)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        f"""---
source_file: "{rel.as_posix()}"
source_type: "{path.suffix.lower().lstrip('.')}"
source_mtime: {current_source_mtime}
generated_from_source: true
content_hash: "{content_hash}"
resolved_path: "{path.resolve().as_posix()}"
mime_type: "{decision.mime_type or 'unknown'}"
extractor_version: "{EXTRACTOR_VERSIONS.get(path.suffix.lower(), 'unknown')}"
embedding_model: "{EMBED_MODEL}"
chunking_version: "{CHUNKING_VERSION}"
indexed_at: "{datetime.now(timezone.utc).isoformat()}"
policy_decision: "{decision.action}"
sensitivity_flags: {json.dumps(decision.sensitivity_flags)}
{extra_frontmatter}---

# {path.stem}

{text.strip()}
""",
        encoding="utf-8",
    )
    return out_path


_chroma_client = None
_collection = None


def get_collection():
    global _chroma_client, _collection
    if _collection is None:
        import chromadb

        CHROMA_DIR.mkdir(parents=True, exist_ok=True)
        _chroma_client = chromadb.PersistentClient(path=str(CHROMA_DIR))
        _collection = _chroma_client.get_or_create_collection(COLLECTION_NAME)
    return _collection
